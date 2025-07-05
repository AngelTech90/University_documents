from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
import os

def create_free_books_ecommerce_presentation():
    """Crear una presentación integral sobre la plataforma e-commerce de libros gratuitos"""
    
    # Create presentation object
    prs = Presentation()
    
    # Define modern color scheme
    PRIMARY_BLUE = RGBColor(41, 128, 185)
    DARK_BLUE = RGBColor(21, 67, 96)
    SUCCESS_GREEN = RGBColor(39, 174, 96)
    WARNING_ORANGE = RGBColor(230, 126, 34)
    DANGER_RED = RGBColor(231, 76, 60)
    PURPLE = RGBColor(142, 68, 173)
    DARK_GRAY = RGBColor(52, 73, 94)
    LIGHT_GRAY = RGBColor(149, 165, 166)
    
    def add_title_slide(title, subtitle=""):
        """Add a professional title slide"""
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        if subtitle:
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(28)
        
        return slide
    
    def add_content_slide(title, content_points, use_icons=False):
        """Add a content slide with enhanced formatting"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        tf = content_shape.text_frame
        tf.clear()
        
        for i, point in enumerate(content_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = point
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(8)
        
        return slide
    
    def add_feature_slide(title, features_dict):
        """Add a slide with feature descriptions"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.color.rgb = DARK_BLUE
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add features in grid layout
        y_pos = 1.5
        for feature, description in features_dict.items():
            # Feature title
            feature_shape = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.4), Inches(0.5))
            feature_frame = feature_shape.text_frame
            feature_frame.text = f"🔹 {feature}"
            feature_frame.paragraphs[0].font.size = Pt(22)
            feature_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
            feature_frame.paragraphs[0].font.bold = True
            
            # Feature description
            desc_shape = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.4), Inches(8), Inches(0.6))
            desc_frame = desc_shape.text_frame
            desc_frame.text = description
            desc_frame.paragraphs[0].font.size = Pt(16)
            desc_frame.paragraphs[0].font.color.rgb = DARK_GRAY
            
            y_pos += 1.2
        
        return slide
    
    def add_comparison_slide(title, left_title, left_points, right_title, right_points):
        """Add a comparison slide with two columns"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.color.rgb = DARK_BLUE
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Left column
        left_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.2), Inches(5.5))
        left_frame = left_shape.text_frame
        left_frame.text = left_title
        left_frame.paragraphs[0].font.size = Pt(24)
        left_frame.paragraphs[0].font.color.rgb = SUCCESS_GREEN
        left_frame.paragraphs[0].font.bold = True
        
        for point in left_points:
            p = left_frame.add_paragraph()
            p.text = f"✓ {point}"
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(6)
        
        # Right column
        right_shape = slide.shapes.add_textbox(Inches(5.3), Inches(1.3), Inches(4.2), Inches(5.5))
        right_frame = right_shape.text_frame
        right_frame.text = right_title
        right_frame.paragraphs[0].font.size = Pt(24)
        right_frame.paragraphs[0].font.color.rgb = WARNING_ORANGE
        right_frame.paragraphs[0].font.bold = True
        
        for point in right_points:
            p = right_frame.add_paragraph()
            p.text = f"⚠ {point}"
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(6)
        
        return slide
    
    def add_technical_slide(title, tech_stack):
        """Add a technical architecture slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.color.rgb = DARK_BLUE
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Create tech stack boxes
        colors = [PRIMARY_BLUE, SUCCESS_GREEN, WARNING_ORANGE, PURPLE]
        x_positions = [1, 3, 5, 7]
        
        for i, (category, technologies) in enumerate(tech_stack.items()):
            # Category box
            box_shape = slide.shapes.add_textbox(Inches(x_positions[i % 4]), Inches(1.8), Inches(1.8), Inches(3.5))
            box_frame = box_shape.text_frame
            box_frame.text = category
            box_frame.paragraphs[0].font.size = Pt(18)
            box_frame.paragraphs[0].font.color.rgb = colors[i % 4]
            box_frame.paragraphs[0].font.bold = True
            box_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Technologies list
            for tech in technologies:
                p = box_frame.add_paragraph()
                p.text = f"• {tech}"
                p.font.size = Pt(14)
                p.font.color.rgb = DARK_GRAY
                p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    # Create slides content
    print("Creando Presentación de Plataforma E-commerce de Libros Gratuitos...")
    
    # Slide 1: Title Slide
    add_title_slide(
        "📚 LibraryLink: Plataforma E-commerce de Libros Gratuitos",
        "Democratizando el Acceso al Conocimiento a través de la Innovación Digital"
    )
    
    # Slide 2: Project Overview
    add_content_slide(
        "🎯 Misión y Visión del Proyecto",
        [
            "Misión: Hacer que el contenido educativo de alta calidad sea accesible gratuitamente para todos",
            "Visión: Crear la plataforma curada más grande del mundo para libros digitales gratuitos",
            "Objetivo: Estudiantes, investigadores, profesionales y aprendices de por vida",
            "Impacto: Cerrar la brecha digital en recursos educativos",
            "Sostenibilidad: Ingresos a través de publicidad ética y patrocinios"
        ]
    )
    
    # Slide 3: Market Problem
    add_content_slide(
        "📊 Análisis del Problema de Mercado",
        [
            "Desigualdad educativa: 60% de estudiantes carecen de acceso a libros de texto de calidad",
            "Altos costos: Los precios de libros de texto han aumentado 812% desde 1978",
            "Fragmentación digital: Libros gratuitos dispersos en múltiples plataformas",
            "Preocupaciones de calidad: Falta de curación y sistemas de reseñas",
            "Desafíos de descubrimiento: Los usuarios luchan por encontrar contenido relevante",
            "Barreras idiomáticas: Recursos educativos multilingües limitados"
        ]
    )
    
    # Slide 4: Solution Overview
    add_content_slide(
        "💡 Nuestra Solución Innovadora",
        [
            "Plataforma centralizada que agrega libros gratuitos de fuentes verificadas",
            "Sistema de curación avanzado con reseñas de expertos y calificaciones",
            "Motor de recomendaciones impulsado por IA para descubrimiento personalizado",
            "Aseguramiento de calidad impulsado por la comunidad a través de comentarios de usuarios",
            "Soporte multiidioma con capacidades de traducción",
            "Diseño accesible siguiendo las pautas WCAG 2.1"
        ]
    )
    
    # Slide 5: Core Features
    add_feature_slide(
        "🚀 Características Principales de la Plataforma",
        {
            "Búsqueda Inteligente y Descubrimiento": "Búsqueda avanzada con filtros por materia, nivel, idioma y formato",
            "Sistema de Reseñas de Expertos": "Bibliotecarios profesionales y educadores proporcionan reseñas detalladas de libros",
            "Sistema de Calificación de 5 Estrellas": "Calificaciones impulsadas por la comunidad ayudan a los usuarios a identificar contenido de calidad",
            "Recomendaciones Personalizadas": "Algoritmos de IA sugieren libros basados en historial de lectura y preferencias",
            "Listas de Lectura y Colecciones": "Los usuarios pueden crear y compartir colecciones curadas de libros",
            "Características Sociales": "Clubes de lectura, foros de discusión y desafíos de lectura"
        }
    )
    
    # Slide 6: Business Model
    add_content_slide(
        "💰 Modelo de Negocio Sostenible",
        [
            "Ingresos Principales: Publicidad contextual relevante al contenido educativo",
            "Contenido Patrocinado: Instituciones educativas y editoriales patrocinan recomendaciones de libros",
            "Características Premium: Analíticas avanzadas y listas de lectura ilimitadas ($5/mes)",
            "Asociaciones Corporativas: Soluciones B2B para escuelas y bibliotecas",
            "Marketing de Afiliados: Asociaciones éticas con servicios educativos",
            "Donaciones: Contribuciones opcionales de usuarios para apoyar el desarrollo de la plataforma"
        ]
    )
    

    
    # Slide 7: User Experience Flow
    add_content_slide(
        "👤 Jornada de Experiencia del Usuario",
        [
            "1. Descubrimiento: Los usuarios llegan a la página de inicio con libros populares y categorías",
            "2. Búsqueda: Búsqueda avanzada con autocompletado inteligente y filtros",
            "3. Evaluación: Ver detalles del libro, reseñas de expertos y calificaciones de la comunidad",
            "4. Acceso: Acceso con un clic a libros gratuitos de fuentes verificadas",
            "5. Participación: Calificar libros, escribir reseñas y unirse a discusiones",
            "6. Personalización: Recibir recomendaciones personalizadas y listas de lectura"
        ]
    )
    
    # Slide 8: Quality Assurance
    add_content_slide(
        "✅ Marco de Aseguramiento de Calidad",
        [
            "Verificación de Fuentes: Solo libros de fuentes legítimas y que cumplen con derechos de autor",
            "Curación de Expertos: Bibliotecarios profesionales revisan todas las adiciones de libros",
            "Moderación Comunitaria: Contenido reportado por usuarios revisado dentro de 24 horas",
            "Escaneo Automatizado: Herramientas de IA detectan contenido potencialmente dañino o inapropiado",
            "Pruebas de Accesibilidad: Todos los libros probados para compatibilidad con lectores de pantalla",
            "Auditorías Regulares: Revisiones mensuales de calidad y detección de enlaces rotos"
        ]
    )
    
    # Slide 9: Monetization Strategy
    add_comparison_slide(
        "💼 Análisis de Fuentes de Ingresos",
        "Fuentes de Ingresos Principales",
        [
            "Publicidad de display (70% de ingresos)",
            "Recomendaciones de libros patrocinadas (20%)",
            "Suscripciones premium (8%)",
            "Asociaciones corporativas (2%)"
        ],
        "Estrategias de Crecimiento",
        [
            "Expandir a instituciones académicas",
            "Desarrollar aplicaciones móviles",
            "Agregar soporte para audiolibros",
            "Expansión a mercados internacionales"
        ]
    )
    
    # Slide 10: Marketing Strategy
    add_content_slide(
        "📈 Estrategia de Marketing y Crecimiento",
        [
            "Marketing de Contenido: Blog educativo con consejos de estudio y recomendaciones de libros",
            "Redes Sociales: Presencia activa en Twitter, LinkedIn y grupos educativos de Facebook",
            "Optimización SEO: Dirigir palabras clave de cola larga relacionadas con recursos educativos gratuitos",
            "Construcción de Comunidad: Asociarse con organizaciones estudiantiles y grupos de estudio en línea",
            "Asociaciones con Influencers: Colaborar con educadores e influencers académicos",
            "Email Marketing: Boletín semanal con recomendaciones curadas de libros"
        ]
    )
    
    # Slide 11: Competitive Analysis
    add_comparison_slide(
        "🏆 Ventaja Competitiva",
        "Nuestras Fortalezas",
        [
            "Contenido curado de calidad",
            "Sistema de reseñas de expertos",
            "Motor de recomendaciones avanzado",
            "Calificaciones impulsadas por la comunidad",
            "Diseño centrado en accesibilidad",
            "Modelo de ingresos éticos"
        ],
        "Diferenciadores de Mercado",
        [
            "Enfoque en contenido educativo",
            "Proceso de curación profesional",
            "Soporte multiidioma",
            "Características de aprendizaje social",
            "Sourcing transparente",
            "Experiencia de lectura sin anuncios"
        ]
    )
    

    
    return prs

def main():
    """Función principal para generar la presentación"""
    print("🚀 Generando Presentación de Plataforma E-commerce de Libros Gratuitos...")
    
    try:
        # Create the presentation
        presentation = create_free_books_ecommerce_presentation()
        
        # Save the presentation
        filename = "LibraryLink_Plataforma_Libros_Gratuitos_Presentacion.pptx"
        presentation.save(filename)
        
        print(f"✅ Presentación guardada exitosamente como '{filename}'")
        print(f"📊 Total de diapositivas creadas: {len(presentation.slides)}")
        print(f"📍 Ubicación del archivo: {os.path.abspath(filename)}")
        
        print("\n📋 Estructura de la Presentación:")
        print("1. Título e Introducción")
        print("2. Misión y Visión del Proyecto")
        print("3. Análisis del Problema de Mercado")
        print("4. Resumen de Solución Innovadora")
        print("5. Características Principales de la Plataforma")
        print("6. Modelo de Negocio Sostenible")
        print("7. Jornada de Experiencia del Usuario")
        print("8. Marco de Aseguramiento de Calidad")
        print("9. Análisis de Fuentes de Ingresos")
        print("10. Estrategia de Marketing y Crecimiento")
        print("11. Ventaja Competitiva")
        
        print("\n🎯 Características de la Presentación:")
        print("• Esquema de colores profesional y tipografía")
        print("• Análisis empresarial integral")
        print("• Enfoque en modelo de negocio principal")
        print("• Enfoque en experiencia del usuario y calidad")
        print("• Análisis de marketing y competencia")
        print("• Optimizada para presentación de 12 minutos")
        
    except ImportError as e:
        print(f"❌ Biblioteca requerida faltante: {str(e)}")
        print("\n🛠️  Instrucciones de Instalación:")
        print("1. Crear entorno virtual:")
        print("   python3 -m venv entorno_plataforma_libros")
        print("2. Activar entorno:")
        print("   source entorno_plataforma_libros/bin/activate  # Linux/Mac")
        print("   entorno_plataforma_libros\\Scripts\\activate     # Windows")
        print("3. Instalar paquete requerido:")
        print("   pip install python-pptx")
        print("4. Ejecutar el script:")
        print("   python3 main.py")
        print("5. Desactivar cuando termine:")
        print("   deactivate")
        
    except Exception as e:
        print(f"❌ Error creando presentación: {str(e)}")
        print("Por favor verifica los permisos de archivo y el espacio disponible en disco.")

if __name__ == "__main__":
    main()
