from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_network_protocols_presentation():
    
    # Create presentation object
    prs = Presentation()
    
    # Define color scheme
    BLUE = RGBColor(41, 128, 185)
    DARK_BLUE = RGBColor(21, 67, 96)
    GREEN = RGBColor(39, 174, 96)
    ORANGE = RGBColor(230, 126, 34)
    
    def add_title_slide(title, subtitle=""):
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        
        if subtitle:
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.color.rgb = BLUE
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
        
        return slide
    
    def add_content_slide(title, content_points):
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        
        tf = content_shape.text_frame
        tf.clear()
        
        for i, point in enumerate(content_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = point
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(51, 51, 51)
        
        return slide
    
    def add_comparison_slide(title, left_title, left_points, right_title, right_points):
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.color.rgb = DARK_BLUE
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Left column
        left_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(5))
        left_frame = left_shape.text_frame
        left_frame.text = left_title
        left_frame.paragraphs[0].font.size = Pt(24)
        left_frame.paragraphs[0].font.color.rgb = BLUE
        left_frame.paragraphs[0].font.bold = True
        
        for point in left_points:
            p = left_frame.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(51, 51, 51)
        
        # Right column
        right_shape = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4), Inches(5))
        right_frame = right_shape.text_frame
        right_frame.text = right_title
        right_frame.paragraphs[0].font.size = Pt(24)
        right_frame.paragraphs[0].font.color.rgb = GREEN
        right_frame.paragraphs[0].font.bold = True
        
        for point in right_points:
            p = right_frame.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(51, 51, 51)
        
        return slide
    
    # Slide 1: Title Slide
    add_title_slide(
        "Protocolos de Red y Diseño de Infraestructura",
        "Guía Completa para la Selección de Equipos de Red e Implementación de Protocolos"
    )
    
    # Slide 2: Agenda
    add_content_slide("Agenda", [
        "Introducción a los Protocolos de Red",
        "Modelo OSI y Pila de Protocolos",
        "Protocolos de Red Fundamentales (TCP/IP, HTTP, DNS, DHCP)",
        "Protocolos Avanzados (BGP, OSPF, STP)",
        "Criterios de Selección de Equipos de Red",
        "Switch vs Router: Cuándo Usar Cada Uno",
        "Principios de Diseño de Redes",
        "Selección de Protocolos para Diferentes Escenarios",
        "Mejores Prácticas y Recomendaciones"
    ])
    
    # Slide 3: What are Network Protocols?
    add_content_slide("¿Qué son los Protocolos de Red?", [
        "Conjunto de reglas que gobiernan la comunicación entre dispositivos de red",
        "Definen cómo se formatean, transmiten y reciben los datos",
        "Garantizan una comunicación confiable y estandarizada",
        "Permiten la interoperabilidad entre diferentes proveedores y sistemas",
        "Forman la base de la infraestructura de redes moderna",
        "Críticos para la seguridad, rendimiento y escalabilidad de la red"
    ])
    
    # Slide 4: OSI Model Overview
    add_content_slide("Modelo OSI y Capas de Protocolos", [
        "Capa 7 (Aplicación): HTTP, HTTPS, FTP, SMTP, DNS",
        "Capa 6 (Presentación): SSL/TLS, cifrado, compresión",
        "Capa 5 (Sesión): NetBIOS, RPC, gestión de sesiones",
        "Capa 4 (Transporte): TCP, UDP, comunicación basada en puertos",
        "Capa 3 (Red): IP, ICMP, protocolos de enrutamiento (OSPF, BGP)",
        "Capa 2 (Enlace de Datos): Ethernet, Wi-Fi, protocolos de conmutación (STP)",
        "Capa 1 (Física): Cables, fibra, transmisión inalámbrica"
    ])
    
    # Slide 5: Core Network Protocols - TCP/IP
    add_content_slide("Suite de Protocolos TCP/IP", [
        "Protocolo de Control de Transmisión (TCP): Confiable, orientado a conexión",
        "Protocolo de Internet (IP): Base de direccionamiento y enrutamiento",
        "Protocolo de Datagrama de Usuario (UDP): Rápido, sin conexión",
        "Protocolo de Mensajes de Control de Internet (ICMP): Informes de errores y diagnósticos",
        "Protocolo de Resolución de Direcciones (ARP): Mapeo de direcciones MAC a IP",
        "IPv4 vs IPv6: Espacio de direcciones y requisitos modernos"
    ])
    
    # Slide 6: Application Layer Protocols
    add_content_slide("Protocolos de Aplicación Esenciales", [
        "HTTP/HTTPS: Comunicación web y transacciones seguras",
        "DNS: Resolución de nombres de dominio y jerarquía",
        "DHCP: Asignación automática de direcciones IP y gestión",
        "SMTP/POP3/IMAP: Transmisión y recuperación de correo electrónico",
        "FTP/SFTP: Transferencia de archivos y operaciones seguras",
        "SNMP: Monitoreo y gestión de redes"
    ])
    
    # Slide 7: Routing Protocols
    add_content_slide("Protocolos de Enrutamiento", [
        "OSPF (Primero la Ruta Más Corta): Convergencia rápida, diseño jerárquico",
        "BGP (Protocolo de Puerta de Enlace Fronterizo): Enrutamiento inter-dominio, backbone de internet",
        "EIGRP: Protocolo propietario de Cisco, híbrido vector-distancia",
        "RIP: Vector-distancia simple, escalabilidad limitada",
        "Consideraciones de enrutamiento estático vs dinámico",
        "Selección de protocolo basada en tamaño de red y requisitos"
    ])
    
    # Slide 8: Switching Protocols
    add_content_slide("Protocolos de Conmutación y VLAN", [
        "Protocolo de Árbol de Expansión (STP): Prevención de bucles en redes conmutadas",
        "VLAN (Red de Área Local Virtual): Segmentación de red y dominios de difusión",
        "Protocolo de Enlace Troncal VLAN (VTP): Sincronización de base de datos VLAN",
        "Agregación de Enlaces (LACP): Agregación de ancho de banda y redundancia",
        "Calidad de Servicio (QoS): Priorización y gestión de tráfico",
        "Seguridad de Puerto: Control de acceso y filtrado de direcciones MAC"
    ])
    
    # Slide 9: Network Equipment Selection Criteria
    add_content_slide("Criterios de Selección de Equipos", [
        "Requisitos de Ancho de Banda: Necesidades de capacidad actuales y futuras",
        "Densidad de Puertos: Número y tipos de conexiones requeridas",
        "Especificaciones de Rendimiento: Throughput, latencia, procesamiento de paquetes",
        "Escalabilidad: Capacidad de crecer con los requisitos del negocio",
        "Redundancia: Alta disponibilidad y capacidades de failover",
        "Características de Gestión: Herramientas de monitoreo, configuración y mantenimiento",
        "Restricciones Presupuestarias: Costo inicial vs valor a largo plazo",
        "Soporte del Proveedor: Documentación, actualizaciones y asistencia técnica"
    ])
    
    # Slide 10: Switch vs Router Comparison
    add_comparison_slide(
        "Switches vs Routers: Diferencias Clave",
        "Switches de Red",
        [
            "Operan en Capa 2 (Enlace de Datos)",
            "Reenvían tramas basándose en direcciones MAC",
            "Crean dominios de colisión",
            "Alta densidad de puertos (típicamente 24-48+ puertos)",
            "Menor latencia, reenvío a velocidad de cable",
            "Soporte VLAN para segmentación",
            "Utilizados dentro de LANs y redes de acceso"
        ],
        "Routers de Red",
        [
            "Operan en Capa 3 (Red)",
            "Enrutan paquetes basándose en direcciones IP",
            "Crean dominios de difusión",
            "Menor densidad de puertos (típicamente 2-16 puertos)",
            "Mayor latencia debido a decisiones de enrutamiento",
            "Soportan múltiples protocolos de enrutamiento",
            "Conectan diferentes redes (LAN-WAN)"
        ]
    )
    
    # Slide 11: When to Use Switches
    add_content_slide("Cuándo Elegir Switches", [
        "Conexiones de capa de acceso de alta densidad (dispositivos de usuario)",
        "Requisitos de conectividad de red de área local (LAN)",
        "Necesidad de segmentación VLAN dentro de la misma subred",
        "Conexiones de granjas de servidores de alto rendimiento",
        "La densidad de puertos rentable es prioridad",
        "Se necesita redundancia de Capa 2 y prevención de bucles",
        "Capas de distribución y acceso de redes de campus"
    ])
    
    # Slide 12: When to Use Routers
    add_content_slide("Cuándo Elegir Routers", [
        "Conectar diferentes redes o subredes",
        "Conectividad WAN y acceso a internet",
        "Requisitos avanzados de protocolos de enrutamiento",
        "Segmentación de red con diferentes rangos IP",
        "Aplicación de políticas de firewall y seguridad",
        "Implementación de Calidad de Servicio (QoS)",
        "Conexiones de internet multi-homed",
        "Conectividad de oficinas sucursales"
    ])
    
    # Slide 13: Network Design Principles
    add_content_slide("Principios de Diseño de Redes", [
        "Diseño Jerárquico: Capas de Núcleo, Distribución y Acceso",
        "Redundancia: Eliminar puntos únicos de falla",
        "Escalabilidad: Diseñar para crecimiento y requisitos futuros",
        "Seguridad: Defensa en profundidad y segmentación",
        "Rendimiento: Planificación de ancho de banda e implementación QoS",
        "Manejabilidad: Monitoreo y configuración centralizados",
        "Rentabilidad: Equilibrar características con restricciones presupuestarias"
    ])
    
    # Slide 14: Protocol Selection Matrix
    add_content_slide("Guías de Selección de Protocolos", [
        "Redes Pequeñas (< 50 dispositivos): Enrutamiento estático, conmutación básica",
        "Redes Medianas (50-500 dispositivos): OSPF, VLANs, QoS básico",
        "Redes Grandes (500+ dispositivos): BGP, enrutamiento avanzado, QoS completo",
        "Requisitos de Alta Seguridad: VPNs, 802.1X, segmentación de red",
        "Aplicaciones en Tiempo Real: QoS, moldeado de tráfico, VLANs dedicadas",
        "Organizaciones Multi-sitio: MPLS, VPNs sitio-a-sitio, gestión centralizada"
    ])
    
    # Slide 15: Common Protocol Combinations
    add_content_slide("Pilas de Protocolos Comunes", [
        "Servicios Web: HTTP/HTTPS + TCP + IP + Ethernet",
        "Sistemas de Correo: SMTP/IMAP + TCP + IP + protocolos de enrutamiento",
        "Compartición de Archivos: FTP/SMB + TCP + IP + protocolos de conmutación",
        "Sistemas VoIP: SIP/RTP + UDP + IP + protocolos QoS",
        "Gestión de Red: SNMP + UDP + IP + protocolos de enrutamiento",
        "Redes Inalámbricas: 802.11 + WPA/WPA2 + IP + protocolos de movilidad"
    ])
    
    # Slide 16: Network Design Process
    add_content_slide("Metodología de Diseño de Redes", [
        "Análisis de Requisitos: Recopilar necesidades técnicas y de negocio",
        "Evaluación de Red: Evaluar infraestructura existente",
        "Diseño de Topología: Crear diseños lógicos y físicos",
        "Selección de Protocolos: Elegir protocolos apropiados para cada capa",
        "Especificación de Equipos: Seleccionar switches, routers y otros dispositivos",
        "Diseño de Seguridad: Implementar estrategias de defensa en profundidad",
        "Planificación de Implementación: Enfoque de despliegue por fases",
        "Pruebas y Validación: Verificar funcionalidad y rendimiento"
    ])
    
    # Slide 17: Performance Considerations
    add_content_slide("Consideraciones de Rendimiento", [
        "Planificación de Ancho de Banda: Calcular requisitos actuales y futuros",
        "Minimización de Latencia: Optimizar decisiones de enrutamiento y conmutación",
        "Implementación QoS: Priorizar flujos de tráfico críticos",
        "Balanceo de Carga: Distribuir tráfico a través de múltiples rutas",
        "Estrategias de Cache: Reducir utilización de ancho de banda WAN",
        "Ajuste de Protocolos: Optimizar ventanas TCP, timeouts y buffers"
    ])
    
    # Slide 18: Security Integration
    add_content_slide("Integración de Protocolos de Seguridad", [
        "Segmentación de Red: VLANs y subredes para aislamiento",
        "Control de Acceso: Autenticación y autorización 802.1X",
        "Cifrado: IPSec, SSL/TLS para protección de datos",
        "Monitoreo: SNMP, syslog y análisis de flujo",
        "Prevención de Intrusiones: Integración IPS con infraestructura de conmutación",
        "Cumplimiento: Requisitos regulatorios y pistas de auditoría"
    ])
    
    # Slide 19: Troubleshooting and Monitoring
    add_content_slide("Monitoreo y Resolución de Problemas", [
        "Analizadores de Protocolos: Wireshark, tcpdump para análisis de paquetes",
        "Monitoreo SNMP: Métricas de estado y rendimiento en tiempo real",
        "Análisis de Flujo: NetFlow, sFlow para patrones de tráfico",
        "Registro y Alertas: Sistemas centralizados de gestión de logs",
        "Líneas Base de Rendimiento: Establecer parámetros normales de operación",
        "Documentación: Diagramas de red, configuraciones y procedimientos"
    ])
    
    # Slide 20: Best Practices Summary
    add_content_slide("Resumen de Mejores Prácticas", [
        "Planificar para crecimiento de 3-5 años al seleccionar equipos",
        "Implementar redundancia en todas las capas críticas de red",
        "Usar protocolos estándar para asegurar interoperabilidad de proveedores",
        "Documentar todas las configuraciones y cambios de red",
        "Actualizaciones regulares de firmware y parches de seguridad",
        "Realizar evaluaciones y optimizaciones periódicas de red",
        "Capacitar al personal en fundamentos de protocolos y resolución de problemas",
        "Mantener relaciones con proveedores para soporte y consultoría"
    ])
    
    # Slide 21: Future Considerations
    add_content_slide("Tecnologías Emergentes y Tendencias", [
        "Redes Definidas por Software (SDN): Planos de control centralizados",
        "Virtualización de Funciones de Red (NFV): Servicios de red virtuales",
        "Redes Basadas en Intención: Implementación automatizada de políticas",
        "Adopción IPv6: Espacio de direcciones y nuevas características de protocolo",
        "Integración IoT: Requisitos de conectividad masiva de dispositivos",
        "Integración en la Nube: Redes híbridas y multi-nube",
        "Redes 5G: Aplicaciones de ultra-baja latencia y alto ancho de banda"
    ])
    
    # Slide 22: Conclusion
    add_content_slide("Puntos Clave", [
        "La selección de protocolos impulsa las decisiones de arquitectura de red",
        "La elección de equipos debe alinearse con los requisitos de protocolos",
        "Diseñar para escalabilidad, seguridad y rendimiento desde el inicio",
        "La evaluación y optimización regular aseguran el éxito continuo",
        "Mantenerse informado sobre tecnologías y estándares emergentes",
        "Invertir en capacitación del equipo y prácticas de documentación"
    ])
    
    return prs

def main():
    print("Generating Network Protocols PowerPoint presentation...")
    
    try:
        # Create the presentation
        presentation = create_network_protocols_presentation()
        
        # Save the presentation
        filename = "./generated_documents/Protocolos_de_Red_y_Diseño_de_Infraestructura.pptx"
        presentation.save(filename)
        
        print(f"✓ Presentación guardada exitosamente como '{filename}'")
        print(f"✓ Total de diapositivas creadas: {len(presentation.slides)}")
        print(f"✓ Ubicación del archivo: {os.path.abspath(filename)}")
        print("\nContenido de la Presentación:")
        print("- Fundamentos de protocolos de red")
        print("- Modelo OSI y capas de protocolos")
        print("- Criterios de selección de equipos")
        print("- Comparación Switch vs Router")
        print("- Principios de diseño de redes")
        print("- Guías de selección de protocolos")
        print("- Mejores prácticas y recomendaciones")
        
    except ImportError as e:
        print(f"✗ Missing required library: {str(e)}")
        print("\nTo fix this, run the following commands:")
        print("1. Create a virtual environment:")
        print("   python3 -m venv network_presentation_env")
        print("2. Activate the virtual environment:")
        print("   source network_presentation_env/bin/activate")
        print("3. Install the required package:")
        print("   pip install python-pptx")
        print("4. Run the script again:")
        print("   python3 main.py")
        print("5. When done, deactivate the environment:")
        print("   deactivate")
        
    except Exception as e:
        print(f"✗ Error creating presentation: {str(e)}")

if __name__ == "__main__":
    main()
